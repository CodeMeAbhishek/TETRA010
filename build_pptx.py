"""
Nidan — Rebuilt Technical Deck (v2)
Tighter, 9-slide deck focused on system design & architecture decisions.
Embeds actual reference screenshots. Plain language, less medical jargon.
Run: python build_pptx.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import os

# ─── PALETTE ──────────────────────────────────────────────────────────────────
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE  = RGBColor(0xF8, 0xFA, 0xFC)
SLATE_900  = RGBColor(0x0F, 0x17, 0x2A)
SLATE_700  = RGBColor(0x33, 0x41, 0x55)
SLATE_500  = RGBColor(0x64, 0x74, 0x8B)
SLATE_200  = RGBColor(0xE2, 0xE8, 0xF0)
TEAL       = RGBColor(0x00, 0x96, 0x88)
TEAL_DARK  = RGBColor(0x00, 0x69, 0x5C)
TEAL_LIGHT = RGBColor(0xB2, 0xDF, 0xDB)
RISK_GREEN  = RGBColor(0x16, 0xA3, 0x4A)
RISK_YELLOW = RGBColor(0xCA, 0x8A, 0x04)
RISK_ORANGE = RGBColor(0xEA, 0x58, 0x0C)
RISK_RED    = RGBColor(0xDC, 0x26, 0x26)
GOLD        = RGBColor(0xF5, 0x9E, 0x0B)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

REF = r"c:\NUV\Tetrathon\references"
IMG_KDIGO_HEATMAP = os.path.join(REF, "kdigo_heatmap (2).png")   # the colour grid
IMG_WHO_LAB       = os.path.join(REF, "south-asia-1.png")
IMG_WHO_NONLAB    = os.path.join(REF, "south-asia-2.png")

# ─── HELPERS ──────────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs

def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def fill_bg(slide, color):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = color

def rect(slide, l, t, w, h, fill=None, line=None, lw=Pt(0)):
    sh = slide.shapes.add_shape(1, l, t, w, h)
    sh.fill.solid() if fill else sh.fill.background()
    if fill: sh.fill.fore_color.rgb = fill
    if line:
        sh.line.color.rgb = line
        sh.line.width = lw
    else:
        sh.line.fill.background()
    return sh

def txt(slide, text, l, t, w, h, size=Pt(14), bold=False, color=SLATE_900,
        align=PP_ALIGN.LEFT, italic=False, name="Calibri"):
    bx = slide.shapes.add_textbox(l, t, w, h)
    tf = bx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size  = size
    r.font.bold  = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = name
    return bx

def badge(slide, text, l, t, w, h, bg=TEAL, fg=WHITE, size=Pt(11), bold=True, align=PP_ALIGN.CENTER):
    sh = rect(slide, l, t, w, h, fill=bg)
    sp = sh._element
    pg = sp.find(qn('p:spPr')).find(qn('a:prstGeom'))
    if pg is not None:
        pg.set('prst', 'roundRect')
        av = pg.find(qn('a:avLst'))
        if av is None: av = etree.SubElement(pg, qn('a:avLst'))
        for g in av.findall(qn('a:gd')): av.remove(g)
        gd = etree.SubElement(av, qn('a:gd'))
        gd.set('name', 'adj'); gd.set('fmla', 'val 18000')
    tf = sh.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = size; r.font.bold = bold
    r.font.color.rgb = fg; r.font.name = "Calibri"
    return sh

def arrow_down(slide, cx, top, h, color=SLATE_500, width=Pt(2)):
    c = slide.shapes.add_connector(1, cx, top, cx, top + h)
    c.line.color.rgb = color
    c.line.width = width
    ln = c._element.find('.//' + qn('a:ln'))
    if ln is not None:
        te = etree.SubElement(ln, qn('a:tailEnd')); te.set('type','none')
        he = etree.SubElement(ln, qn('a:headEnd'))
        he.set('type','arrow'); he.set('w','med'); he.set('len','med')

def header_bar(slide, title, slide_num, dark=True):
    bg = SLATE_900 if dark else RGBColor(0x0B,0x12,0x20)
    rect(slide, 0, 0, SLIDE_W, Inches(1.0), fill=bg)
    txt(slide, title, Inches(0.45), Inches(0.22), Inches(10.5), Inches(0.65),
        size=Pt(28), bold=True, color=WHITE)
    badge(slide, f"  {slide_num}  ", Inches(12.25), Inches(0.28), Inches(0.75), Inches(0.38),
          bg=TEAL, size=Pt(10))

def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text

def add_image(slide, path, l, t, w, h):
    try:
        return slide.shapes.add_picture(path, l, t, w, h)
    except Exception as e:
        # fallback: draw a placeholder box
        r = rect(slide, l, t, w, h, fill=SLATE_200, line=SLATE_500, lw=Pt(1))
        txt(slide, f"[Image: {os.path.basename(path)}]", l+Inches(0.1), t+h/2-Inches(0.2),
            w-Inches(0.2), Inches(0.4), size=Pt(10), color=SLATE_700, align=PP_ALIGN.CENTER)
        return r

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ══════════════════════════════════════════════════════════════════════════════
def slide1(prs):
    sl = blank_slide(prs)
    fill_bg(sl, SLATE_900)
    rect(sl, 0, 0, Inches(0.15), SLIDE_H, fill=TEAL)          # left accent bar
    rect(sl, 0, 0, SLIDE_W, Inches(0.08), fill=TEAL)           # top accent bar

    # Right side coloured condition strip
    cond_cols = [RISK_GREEN, TEAL, RISK_ORANGE, RISK_RED, SLATE_700]
    conds     = ["Diabetes", "Hypertension", "CKD", "CVD", "Stroke"]
    for i,(c,col) in enumerate(zip(conds,cond_cols)):
        badge(sl, c, Inches(9.3), Inches(1.5 + i*1.0), Inches(3.6), Inches(0.62),
              bg=col, size=Pt(14), bold=True)

    badge(sl, "Track A  |  HealthTech  |  Problem Statement 1",
          Inches(0.4), Inches(0.5), Inches(5.8), Inches(0.42), bg=TEAL, size=Pt(11))

    txt(sl, "Nidan", Inches(0.4), Inches(1.2), Inches(8.5), Inches(1.5),
        size=Pt(72), bold=True, color=WHITE)
    txt(sl, "Clinical Decision Support System",
        Inches(0.4), Inches(2.7), Inches(8.5), Inches(0.65),
        size=Pt(24), color=TEAL_LIGHT)
    txt(sl, "AI-powered lifestyle disease risk prediction &\nearly referral assistant for primary care",
        Inches(0.4), Inches(3.45), Inches(8.3), Inches(0.9),
        size=Pt(15), color=RGBColor(0xCB,0xD5,0xE1))

    rect(sl, Inches(0.4), Inches(4.5), Inches(7.5), Inches(0.04), fill=TEAL)

    # Key stats
    stats = [("5", "Conditions covered"), ("4", "Scoring engines"), ("0", "External medical APIs called")]
    for i,(num,lbl) in enumerate(stats):
        sx = Inches(0.5 + i*2.5)
        txt(sl, num, sx, Inches(4.7), Inches(0.8), Inches(0.9),
            size=Pt(36), bold=True, color=TEAL)
        txt(sl, lbl, sx+Inches(0.75), Inches(4.82), Inches(1.6), Inches(0.65),
            size=Pt(11), color=SLATE_500)

    txt(sl, "Team — Tetrathon 2026  |  August 1, 2026  |  github.com/CodeMeAbhishek/TETRA010",
        Inches(0.4), Inches(6.9), Inches(9), Inches(0.4),
        size=Pt(11), color=SLATE_500)
    return sl


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — PROBLEM + APPROACH
# ══════════════════════════════════════════════════════════════════════════════
def slide2(prs):
    sl = blank_slide(prs)
    fill_bg(sl, OFF_WHITE)
    header_bar(sl, "Problem & Approach", "2")

    # Problem side
    rect(sl, Inches(0.35), Inches(1.15), Inches(6.0), Inches(5.85),
         fill=WHITE, line=SLATE_200, lw=Pt(1))
    rect(sl, Inches(0.35), Inches(1.15), Inches(6.0), Inches(0.5), fill=SLATE_900)
    txt(sl, "Problem", Inches(0.55), Inches(1.2), Inches(5.5), Inches(0.4),
        size=Pt(14), bold=True, color=WHITE)

    probs = [
        ("Late diagnosis", "77% of lifestyle disease cases in India are only\ncaught when they have already worsened."),
        ("No decision support", "Primary health workers lack tools to screen\nsystematically — screening is inconsistent."),
        ("Incomplete records", "Rural PHCs often have NO lab results.\nSystem must still give useful output."),
    ]
    for i,(t,b) in enumerate(probs):
        py = Inches(1.85 + i*1.7)
        rect(sl, Inches(0.5), py, Inches(5.7), Inches(1.5),
             fill=OFF_WHITE, line=SLATE_200, lw=Pt(1))
        txt(sl, t, Inches(0.65), py+Inches(0.1), Inches(5.3), Inches(0.38),
            size=Pt(13), bold=True, color=SLATE_900)
        txt(sl, b, Inches(0.65), py+Inches(0.48), Inches(5.3), Inches(0.85),
            size=Pt(12), color=SLATE_700)

    # Approach side
    rect(sl, Inches(7.0), Inches(1.15), Inches(6.0), Inches(5.85),
         fill=WHITE, line=SLATE_200, lw=Pt(1))
    rect(sl, Inches(7.0), Inches(1.15), Inches(6.0), Inches(0.5), fill=TEAL_DARK)
    txt(sl, "Our Approach", Inches(7.2), Inches(1.2), Inches(5.5), Inches(0.4),
        size=Pt(14), bold=True, color=WHITE)

    approaches = [
        (RISK_GREEN, "No external medical APIs",
         "We read the original published research papers and WHO/KDIGO guidelines,\nextracted every formula and threshold, then coded them directly in Python."),
        (TEAL, "Deterministic scoring engines",
         "4 independent engines run the math. Each engine cites its source. The AI layer\nonly explains results — it never calculates a score."),
        (RISK_ORANGE, "Graceful degradation",
         "If lab data is missing, the system automatically uses the non-lab fallback chart\nand tells the doctor exactly which test to order — never guesses."),
    ]
    for i,(col,t,b) in enumerate(approaches):
        ay = Inches(1.85 + i*1.7)
        rect(sl, Inches(7.15), ay, Inches(0.12), Inches(1.5), fill=col)
        txt(sl, t, Inches(7.4), ay+Inches(0.1), Inches(5.3), Inches(0.38),
            size=Pt(13), bold=True, color=SLATE_900)
        txt(sl, b, Inches(7.4), ay+Inches(0.48), Inches(5.3), Inches(0.9),
            size=Pt(11.5), color=SLATE_700)

    notes(sl, "KEY MESSAGE: We didn't call any medical API or let an AI estimate a risk number. "
              "We went to the source papers, pulled the validated formulas, and hard-coded them. "
              "That is what makes this auditable.")
    return sl


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — SYSTEM ARCHITECTURE
# ══════════════════════════════════════════════════════════════════════════════
def slide3(prs):
    sl = blank_slide(prs)
    fill_bg(sl, SLATE_900)
    header_bar(sl, "System Architecture", "3", dark=False)

    # ── Node layout ──────────────────────────────────────────────────
    cx = Inches(6.65)
    nw = Inches(5.2)
    nh = Inches(0.52)
    gap = Inches(0.32)

    def node(text, top, fill, fg=WHITE, size=Pt(12)):
        badge(sl, text, cx - nw/2, top, nw, nh, bg=fill, fg=fg, size=size)

    def arr(top):
        arrow_down(sl, cx, top, gap, color=TEAL, width=Pt(2))

    tops = [Inches(1.1), Inches(1.94), Inches(2.78), Inches(3.75), Inches(4.75), Inches(5.55), Inches(6.3)]

    node("Frontend  —  Patient Assessment Form (HTML/JS)", tops[0], TEAL_DARK)
    arr(tops[0]+nh)
    node("FastAPI Backend  —  /analyze  endpoint (Python)", tops[1], SLATE_700)
    arr(tops[1]+nh)
    node("Input Checker  —  identifies which fields are present vs. missing", tops[2], RGBColor(0x37,0x47,0x5A))

    # 4 parallel engines
    ew = Inches(2.85)
    eg = Inches(0.15)
    etotal = 4*ew + 3*eg
    ex0 = (SLIDE_W - etotal)/2
    engines = [
        ("Diabetes\nIDRS formula", RISK_GREEN),
        ("Blood Pressure\nACC/AHA 2017", TEAL),
        ("Heart Risk\nWHO CVD charts", RISK_ORANGE),
        ("Kidney Risk\nKDIGO 2024", RISK_RED),
    ]
    engine_y = tops[3]
    for i,(lbl,col) in enumerate(engines):
        bx = ex0 + i*(ew+eg)
        badge(sl, lbl, bx, engine_y, ew, Inches(0.68), bg=col, size=Pt(10.5))
        arrow_down(sl, bx+ew/2, tops[2]+nh, engine_y-(tops[2]+nh), color=TEAL, width=Pt(1.5))

    # converge line + arrow
    conv_y = engine_y + Inches(0.68) + Inches(0.1)
    rect(sl, ex0, conv_y, etotal, Inches(0.03), fill=TEAL)
    arr(conv_y)

    node("Missing Investigations  —  de-duplicates what each engine needs", tops[4], SLATE_700)
    arr(tops[4]+nh)
    node("Referral Engine  —  REFER / MONITOR / ROUTINE + reason", tops[5], SLATE_700)

    # Gold arrow — the key claim
    gold_top = tops[5]+nh+Inches(0.04)
    arrow_down(sl, cx, gold_top, Inches(0.22), color=GOLD, width=Pt(3))
    lbw = Inches(6.4)
    badge(sl, "Structured JSON only passed to AI  —  AI never calculates a score",
          cx-lbw/2, gold_top-Inches(0.02), lbw, Inches(0.36),
          bg=GOLD, fg=SLATE_900, size=Pt(12), bold=True)

    node("AI Layer (NVIDIA NIM / Claude)  —  explain  |  translate  |  draft referral note",
         tops[6], RGBColor(0x31,0x2E,0x81), size=Pt(11))

    notes(sl, "Walk top to bottom. Emphasize:\n"
              "1) All 4 engines run in parallel — no blocking.\n"
              "2) The GOLD arrow: engines produce JSON, AI receives JSON. AI is an explainer, not a calculator.\n"
              "3) Missing investigations aggregator means the doctor gets ONE list, not four.")
    return sl


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — HOW WE BUILT EACH ENGINE (no API — straight from papers)
# ══════════════════════════════════════════════════════════════════════════════
def slide4(prs):
    sl = blank_slide(prs)
    fill_bg(sl, OFF_WHITE)
    header_bar(sl, "How We Built the Scoring Engines  —  No Medical API", "4")

    # Callout box
    rect(sl, Inches(0.35), Inches(1.1), SLIDE_W-Inches(0.7), Inches(0.65), fill=GOLD)
    txt(sl, "Every formula in this system came from reading the original published research papers, "
            "not from calling any third-party medical API.",
        Inches(0.55), Inches(1.18), SLIDE_W-Inches(1.1), Inches(0.5),
        size=Pt(13), bold=True, color=SLATE_900, align=PP_ALIGN.CENTER)

    engines = [
        {
            "name": "Diabetes Risk  (IDRS)",
            "color": RISK_GREEN,
            "source": "ICMR-INDIAB study  |  113,043 patients across India",
            "what": "We read the research paper and extracted a 4-factor scoring formula:\nAge + Waist size + Physical activity + Family history  =  score 0 to 100",
            "code": "if waist_cm >= 90: score += 30   # male threshold\nif age >= 50:       score += 20\nif sedentary:       score += 20\nif family_history:  score += 10",
        },
        {
            "name": "Heart Risk  (WHO CVD Charts)",
            "color": RISK_ORANGE,
            "source": "WHO South Asia CVD Risk Charts  —  digitised into JSON grid",
            "what": "We photographed and manually digitised the WHO colour risk table\n(800+ cells). System looks up the right cell for each patient.",
            "code": "# No API — we digitised the chart ourselves\nrisk = who_cvd_grid[age_band][sex][sbp_band][cholesterol_band]",
        },
        {
            "name": "Kidney Risk  (KDIGO)",
            "color": RISK_RED,
            "source": "KDIGO 2024 Clinical Practice Guidelines  +  CKD-EPI 2021 formula",
            "what": "We extracted the eGFR formula from the guideline and coded the\ncolour risk grid (G1-G5 x A1-A3 = 15 cells) directly in Python.",
            "code": "eGFR = 142 * min(cr/kappa,1)**alpha * max(cr/kappa,1)**-1.200\nrisk_color = kdigo_grid[g_stage][a_stage]",
        },
        {
            "name": "Blood Pressure  (ACC/AHA 2017)",
            "color": TEAL,
            "source": "2017 ACC/AHA Hypertension Guideline  +  Stroke protocols",
            "what": "We read the guideline document and coded every BP threshold,\ncross-module rule (e.g. CKD+HTN = stricter target), and stroke protocol.",
            "code": "if sbp >= 180 or dbp >= 110: category = 'Stage 2 Hypertension'\nelif sbp >= 130:             category = 'Stage 1 Hypertension'",
        },
    ]

    # 2x2 grid
    for i, eng in enumerate(engines):
        col_i = i % 2
        row_i = i // 2
        bx = Inches(0.35 + col_i*6.55)
        by = Inches(2.0 + row_i*2.55)
        rect(sl, bx, by, Inches(6.3), Inches(2.4), fill=WHITE, line=eng["color"], lw=Pt(1.5))
        rect(sl, bx, by, Inches(6.3), Inches(0.48), fill=eng["color"])
        txt(sl, eng["name"], bx+Inches(0.15), by+Inches(0.07), Inches(5.9), Inches(0.35),
            size=Pt(13), bold=True, color=WHITE)
        txt(sl, "Source: "+eng["source"], bx+Inches(0.15), by+Inches(0.55),
            Inches(5.9), Inches(0.32), size=Pt(10), color=TEAL_DARK, italic=True)
        txt(sl, eng["what"], bx+Inches(0.15), by+Inches(0.88),
            Inches(5.9), Inches(0.65), size=Pt(10.5), color=SLATE_900)
        rect(sl, bx+Inches(0.15), by+Inches(1.55), Inches(5.9), Inches(0.75),
             fill=SLATE_900)
        txt(sl, eng["code"], bx+Inches(0.25), by+Inches(1.58),
            Inches(5.7), Inches(0.7), size=Pt(8.5), color=TEAL_LIGHT, name="Courier New")

    notes(sl, "Key talking point: We did the hard manual work of reading these papers and "
              "coding their formulas ourselves. That is WHY the system is accurate and auditable — "
              "there is no black box API that we're hoping is correct.\n"
              "The code snippets shown are simplified but accurate excerpts from the actual engine files.")
    return sl


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — REFERENCE CHARTS (actual images embedded)
# ══════════════════════════════════════════════════════════════════════════════
def slide5(prs):
    sl = blank_slide(prs)
    fill_bg(sl, SLATE_900)
    header_bar(sl, "The Source Charts We Digitised Into Code", "5", dark=False)

    # KDIGO heatmap (left)
    rect(sl, Inches(0.35), Inches(1.1), Inches(5.8), Inches(5.85),
         fill=RGBColor(0x1A,0x28,0x3E))
    txt(sl, "KDIGO Kidney Risk Grid", Inches(0.5), Inches(1.18),
        Inches(5.5), Inches(0.4), size=Pt(13), bold=True, color=WHITE)
    txt(sl, "G1-G5 (kidney function) x A1-A3 (protein in urine)  =  risk colour",
        Inches(0.5), Inches(1.55), Inches(5.5), Inches(0.35),
        size=Pt(10.5), color=TEAL_LIGHT)
    add_image(sl, IMG_KDIGO_HEATMAP, Inches(0.45), Inches(1.95), Inches(5.7), Inches(4.7))
    txt(sl, "We coded all 15 cells of this grid into Python dict — no API",
        Inches(0.5), Inches(6.75), Inches(5.5), Inches(0.35),
        size=Pt(10), color=GOLD, italic=True)

    # WHO lab chart (right top)
    rect(sl, Inches(6.55), Inches(1.1), Inches(6.4), Inches(2.8),
         fill=RGBColor(0x1A,0x28,0x3E))
    txt(sl, "WHO CVD Chart  —  Lab Version (with Cholesterol)",
        Inches(6.7), Inches(1.18), Inches(6.1), Inches(0.35),
        size=Pt(12), bold=True, color=WHITE)
    txt(sl, "Rows: age  |  Columns: cholesterol + BP  |  Colour: 10-year risk %",
        Inches(6.7), Inches(1.52), Inches(6.1), Inches(0.28),
        size=Pt(9.5), color=TEAL_LIGHT)
    add_image(sl, IMG_WHO_LAB, Inches(6.6), Inches(1.85), Inches(6.25), Inches(1.85))

    # WHO non-lab chart (right bottom)
    rect(sl, Inches(6.55), Inches(4.05), Inches(6.4), Inches(2.85),
         fill=RGBColor(0x1A,0x28,0x3E))
    txt(sl, "WHO CVD Chart  —  Non-Lab Version (uses BMI instead)",
        Inches(6.7), Inches(4.13), Inches(6.1), Inches(0.35),
        size=Pt(12), bold=True, color=WHITE)
    txt(sl, "Auto-used when cholesterol lab result is unavailable",
        Inches(6.7), Inches(4.47), Inches(6.1), Inches(0.28),
        size=Pt(9.5), color=RISK_ORANGE)
    add_image(sl, IMG_WHO_NONLAB, Inches(6.6), Inches(4.8), Inches(6.25), Inches(1.95))

    txt(sl, "These are the ACTUAL source images. We manually digitised them — "
            "every number in our JSON data files came from these charts.",
        Inches(0.35), Inches(6.78), SLIDE_W-Inches(0.7), Inches(0.38),
        size=Pt(11), color=GOLD, align=PP_ALIGN.CENTER)

    notes(sl, "This slide is the visual proof of the manual work.\n"
              "Point to the KDIGO grid: 'We coded every one of these 15 coloured cells into a Python dict.'\n"
              "Point to WHO lab chart: 'Each of these coloured numbers is a JSON value that our engine looks up.'\n"
              "Point to WHO non-lab: 'This is the fallback chart — uses BMI instead of cholesterol. "
              "Our system automatically switches to this when cholesterol is missing.'")
    return sl


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — GRACEFUL DEGRADATION (missing data handling)
# ══════════════════════════════════════════════════════════════════════════════
def slide6(prs):
    sl = blank_slide(prs)
    fill_bg(sl, OFF_WHITE)
    header_bar(sl, "Handling Missing Data  —  Graceful Degradation", "6")

    txt(sl, "The problem statement explicitly calls out incomplete patient records in rural settings. "
            "Here is exactly how we handle it:",
        Inches(0.4), Inches(1.12), Inches(12.5), Inches(0.42),
        size=Pt(13), color=SLATE_700, align=PP_ALIGN.CENTER)

    # Flow diagram — same patient, two paths
    # Left column header
    badge(sl, "PATIENT ARRIVES",
          Inches(5.2), Inches(1.72), Inches(2.9), Inches(0.48),
          bg=SLATE_700, size=Pt(13))
    arrow_down(sl, Inches(6.65), Inches(2.2), Inches(0.3), color=SLATE_700, width=Pt(2))

    badge(sl, "Are lab results available?",
          Inches(5.0), Inches(2.5), Inches(3.3), Inches(0.44),
          bg=SLATE_900, size=Pt(12))

    # Yes branch
    for step_y, label, col in [
        (Inches(3.3),  "YES — Cholesterol available", RISK_GREEN),
        (Inches(4.05), "Use lab-based WHO-CVD chart", RISK_GREEN),
        (Inches(4.8),  "High-precision 10-year risk %", RISK_GREEN),
    ]:
        badge(sl, label, Inches(1.0), step_y, Inches(4.3), Inches(0.44), bg=col, size=Pt(12))
        if step_y > Inches(3.3):
            arrow_down(sl, Inches(3.15), step_y - Inches(0.3), Inches(0.3), color=col, width=Pt(2))

    # No branch
    for step_y, label, col in [
        (Inches(3.3),  "NO — Cholesterol missing", RISK_ORANGE),
        (Inches(4.05), "Auto-switch to BMI-based chart", RISK_ORANGE),
        (Inches(4.8),  "Risk calculated via BMI instead", RISK_ORANGE),
    ]:
        badge(sl, label, Inches(8.0), step_y, Inches(4.3), Inches(0.44), bg=col, size=Pt(12))
        if step_y > Inches(3.3):
            arrow_down(sl, Inches(10.15), step_y - Inches(0.3), Inches(0.3), color=col, width=Pt(2))

    # Flag box
    rect(sl, Inches(6.4), Inches(3.9), Inches(0.5), Inches(0.5))   # spacer
    rect(sl, Inches(7.9), Inches(5.4), Inches(4.5), Inches(0.82),
         fill=RGBColor(0xFF,0xED,0xCC), line=GOLD, lw=Pt(1.5))
    txt(sl, "Flag raised for doctor:\n\"Lipid panel recommended  —  Source: WHO CVD Guidelines\"",
        Inches(8.05), Inches(5.45), Inches(4.25), Inches(0.72),
        size=Pt(11), color=SLATE_900)

    # CKD missing path
    rect(sl, Inches(0.35), Inches(6.1), SLIDE_W-Inches(0.7), Inches(0.98),
         fill=SLATE_900)
    txt(sl, "Same logic for kidney tests: if Serum Creatinine is missing "
            "-> engine returns partial result + flags \"Serum Creatinine required for kidney risk calculation\"",
        Inches(0.55), Inches(6.22), SLIDE_W-Inches(1.1), Inches(0.72),
        size=Pt(12), color=WHITE)

    notes(sl, "Emphasize: we never silently substitute a default value for missing data. "
              "The system tells the doctor exactly what is missing AND why it matters. "
              "The missing investigations list is de-duplicated — if two engines need creatinine, "
              "the doctor only sees one 'order creatinine' recommendation.")
    return sl


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — TECH STACK & PROGRESS
# ══════════════════════════════════════════════════════════════════════════════
def slide7(prs):
    sl = blank_slide(prs)
    fill_bg(sl, SLATE_900)
    header_bar(sl, "Tech Stack & Progress", "7", dark=False)

    stack_items = [
        ("FastAPI  (Python)", "Backend REST API  —  /analyze endpoint with Pydantic input validation", TEAL),
        ("Pure Python engines", "4 scoring engines — zero ML/AI dependencies for the math. 86% unit test coverage", RISK_GREEN),
        ("NVIDIA NIM  /  Claude AI", "AI layer: generates the plain-English explanation and referral note only.\nReceives structured JSON — never computes a risk score", GOLD),
        ("Vanilla HTML / CSS / JS", "Frontend — no framework. Runs in any browser, works on a tablet.\n10-patient demo dropdown built in for live evaluation", SLATE_500),
        ("Synthetic test suite", "10 clinical edge-case patients (missing labs, elderly, stroke, healthy baseline)\nCaught 3 real bugs before demo", RISK_ORANGE),
    ]

    for i,(tech,detail,col) in enumerate(stack_items):
        ry = Inches(1.18 + i*1.15)
        rect(sl, Inches(0.35), ry, Inches(12.6), Inches(1.05),
             fill=RGBColor(0x1A,0x28,0x3E) if i%2==0 else RGBColor(0x15,0x20,0x33))
        rect(sl, Inches(0.35), ry, Inches(0.12), Inches(1.05), fill=col)
        badge(sl, tech, Inches(0.6), ry+Inches(0.28), Inches(2.8), Inches(0.42),
              bg=col, size=Pt(11), fg=WHITE if col!=GOLD else SLATE_900)
        txt(sl, detail, Inches(3.6), ry+Inches(0.12), Inches(9.1), Inches(0.82),
            size=Pt(12), color=RGBColor(0xCB,0xD5,0xE1))

    notes(sl, "Be honest about test coverage: 86%. Stroke edge cases are the thinnest area.\n"
              "AI layer is a graceful enhancement: if the API key is absent, all 4 scoring engines "
              "still run and return full structured data — the AI explanation is additive, not required.")
    return sl


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — BUGS WE CAUGHT & FIXED (credibility)
# ══════════════════════════════════════════════════════════════════════════════
def slide8(prs):
    sl = blank_slide(prs)
    fill_bg(sl, OFF_WHITE)
    header_bar(sl, "What We Caught & Fixed  —  Validation Approach", "8")

    txt(sl, "Honest reporting of real bugs found during development builds more credibility than claiming a perfect build.",
        Inches(0.4), Inches(1.1), Inches(12.5), Inches(0.38),
        size=Pt(12), italic=True, color=SLATE_700, align=PP_ALIGN.CENTER)

    bugs = [
        ("WHO CVD Chart — BMI boundary gap",
         "Our digitised grid had 30–34 where the source chart shows 30–35. "
         "A patient with BMI exactly 34.8 was reading the wrong risk band.",
         "Found by: manually comparing our JSON grid cell-by-cell against the source chart image",
         RISK_ORANGE),
        ("Kidney Risk Grid — wrong cell colour",
         "One specific cell (G4 + low protein) was coded as Yellow but the official KDIGO grid shows Orange.",
         "Found by: comparing our hardcoded grid against the actual KDIGO heatmap screenshot we had saved",
         RISK_RED),
        ("Kidney eGFR — float vs integer boundary",
         "eGFR of 44.2 was falling through our G-stage conditions written as integers (30 <= eGFR <= 44), "
         "giving wrong G5 instead of G3b.",
         "Found by: synthetic test Patient 10 (elderly boundary case) in our automated test suite",
         RISK_RED),
    ]

    for i,(title,desc,found,col) in enumerate(bugs):
        by = Inches(1.75 + i*1.68)
        rect(sl, Inches(0.35), by, Inches(12.6), Inches(1.55),
             fill=WHITE, line=col, lw=Pt(1.5))
        rect(sl, Inches(0.35), by, Inches(0.35), Inches(1.55), fill=col)
        badge(sl, f"Bug {i+1}", Inches(0.85), by+Inches(0.12), Inches(0.9), Inches(0.34),
              bg=col, size=Pt(10))
        txt(sl, title, Inches(1.9), by+Inches(0.12), Inches(10.8), Inches(0.38),
            size=Pt(13.5), bold=True, color=SLATE_900)
        txt(sl, desc, Inches(0.85), by+Inches(0.55), Inches(11.8), Inches(0.42),
            size=Pt(11.5), color=SLATE_700)
        txt(sl, f"  {found}", Inches(0.85), by+Inches(1.05), Inches(11.8), Inches(0.38),
            size=Pt(10.5), italic=True, color=TEAL_DARK)

    rect(sl, Inches(0.35), Inches(6.82), Inches(12.6), Inches(0.52), fill=SLATE_900)
    txt(sl, "10-patient synthetic test suite — covering missing data, boundary values, elderly patients, "
            "acute stroke — caught all 3 bugs before demo.",
        Inches(0.55), Inches(6.88), Inches(12.1), Inches(0.42),
        size=Pt(12), bold=True, color=GOLD, align=PP_ALIGN.CENTER)

    notes(sl, "Don't be defensive about these bugs — they show the validation process is working.\n"
              "Bug 2 was only catchable because we HAD the source image saved locally and could compare.\n"
              "Bug 3 was caught by the test suite, not manually — shows the test suite was worth building.")
    return sl


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — LIVE DEMO + FAQ
# ══════════════════════════════════════════════════════════════════════════════
def slide9(prs):
    sl = blank_slide(prs)
    fill_bg(sl, SLATE_900)
    rect(sl, 0, 0, SLIDE_W, Inches(0.06), fill=TEAL)

    # Left: Demo cue
    rect(sl, 0, 0, Inches(6.5), SLIDE_H, fill=RGBColor(0x07,0x0F,0x1E))
    txt(sl, "Live Demo", Inches(0.4), Inches(2.2), Inches(5.7), Inches(1.4),
        size=Pt(60), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(sl, "localhost:8081", Inches(0.4), Inches(3.65), Inches(5.7), Inches(0.6),
        size=Pt(18), color=TEAL_LIGHT, align=PP_ALIGN.CENTER)

    scenarios = [
        ("P03", "High Risk case — full lab data", RISK_RED),
        ("P04", "Missing labs — graceful fallback", RISK_ORANGE),
        ("P06", "Healthy patient — all green", RISK_GREEN),
    ]
    for i,(pid,lbl,col) in enumerate(scenarios):
        sy = Inches(4.55 + i*0.75)
        badge(sl, pid, Inches(0.55), sy, Inches(0.65), Inches(0.42), bg=col, size=Pt(10))
        txt(sl, lbl, Inches(1.35), sy+Inches(0.06), Inches(4.7), Inches(0.35),
            size=Pt(12.5), color=WHITE)

    # Right: Quick FAQ
    rect(sl, Inches(6.7), 0, Inches(6.63), SLIDE_H, fill=RGBColor(0x0F,0x17,0x2A))
    badge(sl, "Likely Questions", Inches(6.85), Inches(0.3), Inches(5.8), Inches(0.45),
          bg=TEAL, size=Pt(12))

    faqs = [
        ("Why not just ask an AI for the risk score?",
         "AI can hallucinate numbers. Our engines use the actual validated formulas — "
         "every score is traceable to a specific guideline."),
        ("Is this validated on real patients?",
         "The formulas we coded ARE validated — IDRS on 113k patients, WHO/KDIGO by global health bodies. "
         "Our implementation tested on 10 synthetic edge-case patients."),
        ("What if two conditions give conflicting advice?",
         "The ACC/AHA guideline has cross-condition rules (e.g. CKD + HTN = stricter BP target). "
         "We coded these combined rules into the referral engine."),
        ("What would you build next?",
         "Offline mode (PS2 scope), more regional language support, "
         "expanded synthetic test coverage to 100+ cases."),
    ]
    for i,(q,a) in enumerate(faqs):
        fy = Inches(1.0 + i*1.55)
        rect(sl, Inches(6.85), fy, Inches(6.1), Inches(1.42),
             fill=RGBColor(0x1A,0x28,0x3E), line=TEAL, lw=Pt(0.5))
        txt(sl, "Q: "+q, Inches(7.0), fy+Inches(0.1), Inches(5.7), Inches(0.42),
            size=Pt(11), bold=True, color=WHITE)
        txt(sl, a, Inches(7.0), fy+Inches(0.56), Inches(5.7), Inches(0.78),
            size=Pt(10.5), color=RGBColor(0xCB,0xD5,0xE1))

    txt(sl, "github.com/CodeMeAbhishek/TETRA010",
        Inches(6.85), Inches(7.1), Inches(6.1), Inches(0.32),
        size=Pt(11), color=TEAL_LIGHT, align=PP_ALIGN.CENTER)

    notes(sl, "Demo order:\n"
              "1. P03 — walk through all 4 gauges + referral banner\n"
              "2. P04 — show WHO chart fallback + missing investigations panel\n"
              "3. P06 — show calm low-risk state (system doesn't just alarm on everyone)\n\n"
              "After demo, offer to open GitHub and show /backend/scoring/ folder.")
    return sl


# ══════════════════════════════════════════════════════════════════════════════
# MAIN
# ══════════════════════════════════════════════════════════════════════════════
def main():
    prs = new_prs()
    slide1(prs)
    slide2(prs)
    slide3(prs)
    slide4(prs)
    slide5(prs)
    slide6(prs)
    slide7(prs)
    slide8(prs)
    slide9(prs)

    out = r"c:\NUV\Tetrathon\Nidan_TechDeck.pptx"
    prs.save(out)
    print(f"[OK] Saved: {out}")
    print(f"     Slides: {len(prs.slides)}")

if __name__ == "__main__":
    main()
